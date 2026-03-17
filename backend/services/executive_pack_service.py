"""Document-oriented executive PPTX export pipeline for QA Platform."""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

import pandas as pd

from backend.config import get_settings
from backend.services.dossier_service import build_executive_report_payload
from backend.services.piece_signal_service import load_piece_signal_payload

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@dataclass(frozen=True)
class _Theme:
    navy: tuple[int, int, int] = (11, 45, 78)
    slate: tuple[int, int, int] = (79, 95, 115)
    light_bg: tuple[int, int, int] = (246, 248, 250)
    white: tuple[int, int, int] = (255, 255, 255)
    accent: tuple[int, int, int] = (0, 133, 119)
    warning: tuple[int, int, int] = (185, 28, 28)


CONTENT_BOX = {"x": 0.62, "y": 1.62, "w": 10.90, "h": 5.55}
CONTENT_EDGE_PAD = 0.04


def _box(x: float, y: float, w: float, h: float) -> dict[str, float]:
    return {"x": float(x), "y": float(y), "w": float(w), "h": float(h)}


def assertInsideContentBox(box: dict[str, float], contentBox: dict[str, float]) -> None:
    eps = 1e-6
    right = box["x"] + box["w"]
    bottom = box["y"] + box["h"]
    c_right = contentBox["x"] + contentBox["w"]
    c_bottom = contentBox["y"] + contentBox["h"]
    if box["x"] < contentBox["x"] - eps or box["y"] < contentBox["y"] - eps or right > c_right + eps or bottom > c_bottom + eps:
        raise ValueError(f"Shape outside content box: box={box}, contentBox={contentBox}")


def _clamp_to_content_box(box: dict[str, float], content_box: dict[str, float]) -> dict[str, float]:
    inner_x = content_box["x"] + CONTENT_EDGE_PAD
    inner_y = content_box["y"] + CONTENT_EDGE_PAD
    inner_w = max(0.25, content_box["w"] - 2 * CONTENT_EDGE_PAD)
    inner_h = max(0.18, content_box["h"] - 2 * CONTENT_EDGE_PAD)

    max_w = max(0.25, inner_w)
    max_h = max(0.18, inner_h)
    w = min(max(0.25, box["w"]), max_w)
    h = min(max(0.18, box["h"]), max_h)
    max_x = inner_x + inner_w - w
    max_y = inner_y + inner_h - h
    x = min(max(box["x"], inner_x), max_x)
    y = min(max(box["y"], inner_y), max_y)
    clamped = _box(x, y, w, h)
    assertInsideContentBox(clamped, content_box)
    return clamped


def _reflow_inside_content_box(box: dict[str, float], content_box: dict[str, float]) -> dict[str, float]:
    clamped = _clamp_to_content_box(box, content_box)
    try:
        assertInsideContentBox(clamped, content_box)
        return clamped
    except ValueError:
        fallback = _clamp_to_content_box(
            _box(clamped["x"], clamped["y"], max(0.25, clamped["w"] - 0.06), max(0.18, clamped["h"] - 0.06)),
            content_box,
        )
        assertInsideContentBox(fallback, content_box)
        return fallback


def _layout_overview(content_box: dict[str, float]) -> dict[str, Any]:
    pad = CONTENT_EDGE_PAD
    gap = 0.10
    cards_h = 1.00
    inner_x = content_box["x"] + pad
    inner_y = content_box["y"] + pad
    inner_w = content_box["w"] - 2 * pad
    inner_h = content_box["h"] - 2 * pad
    
    cards_w = (inner_w - (3 * gap)) / 4
    cards = []
    for idx in range(4):
        card_x = inner_x + idx * (cards_w + gap)
        card_box = _box(card_x, inner_y, cards_w, cards_h)
        clamped = _reflow_inside_content_box(card_box, content_box)
        assertInsideContentBox(clamped, content_box)
        cards.append(clamped)

    tables_top = inner_y + cards_h + 0.18
    tables_h = inner_h - cards_h - 0.18
    table_w = (inner_w - gap) / 2
    left_table_box = _box(inner_x, tables_top, table_w, tables_h)
    left_table = _reflow_inside_content_box(left_table_box, content_box)
    assertInsideContentBox(left_table, content_box)
    
    right_table_box = _box(inner_x + table_w + gap, tables_top, table_w, tables_h)
    right_table = _reflow_inside_content_box(right_table_box, content_box)
    assertInsideContentBox(right_table, content_box)
    
    return {"cards": cards, "left_table": left_table, "right_table": right_table}


def _layout_trend(content_box: dict[str, float], *, include_secondary_chart: bool) -> dict[str, dict[str, float]]:
    pad = CONTENT_EDGE_PAD
    gap = 0.16
    top_h = 3.35
    
    inner_x = content_box["x"] + pad
    inner_y = content_box["y"] + pad
    inner_w = content_box["w"] - 2 * pad
    inner_h = content_box["h"] - 2 * pad
    bottom_h = inner_h - top_h - gap
    
    if include_secondary_chart:
        secondary_w = max(2.35, min(2.96, inner_w * 0.24))
        main_w = max(5.0, inner_w - secondary_w - gap)
        main_chart_box = _box(inner_x, inner_y, main_w, top_h)
        main_chart = _reflow_inside_content_box(main_chart_box, content_box)
        assertInsideContentBox(main_chart, content_box)
        
        secondary_chart_box = _box(inner_x + main_w + gap, inner_y, secondary_w, top_h)
        secondary_chart = _reflow_inside_content_box(secondary_chart_box, content_box)
        assertInsideContentBox(secondary_chart, content_box)
    else:
        main_chart_box = _box(inner_x, inner_y, inner_w, top_h)
        main_chart = _reflow_inside_content_box(main_chart_box, content_box)
        assertInsideContentBox(main_chart, content_box)
        secondary_chart = _reflow_inside_content_box(_box(inner_x, inner_y, 0.25, 0.18), content_box)

    bottom_table_box = _box(inner_x, inner_y + top_h + gap, inner_w, bottom_h)
    bottom_table = _reflow_inside_content_box(bottom_table_box, content_box)
    assertInsideContentBox(bottom_table, content_box)
    
    return {"main_chart": main_chart, "secondary_chart": secondary_chart, "bottom_table": bottom_table}


def _layout_risk(content_box: dict[str, float]) -> dict[str, dict[str, float]]:
    pad = CONTENT_EDGE_PAD
    gap = 0.10
    top_h = 2.30
    
    inner_x = content_box["x"] + pad
    inner_y = content_box["y"] + pad
    inner_w = content_box["w"] - 2 * pad
    inner_h = content_box["h"] - 2 * pad
    bottom_h = inner_h - top_h - gap
    
    top_w = (inner_w - gap) / 2
    top_left_box = _box(inner_x, inner_y, top_w, top_h)
    top_left = _reflow_inside_content_box(top_left_box, content_box)
    assertInsideContentBox(top_left, content_box)
    
    top_right_box = _box(inner_x + top_w + gap, inner_y, top_w, top_h)
    top_right = _reflow_inside_content_box(top_right_box, content_box)
    assertInsideContentBox(top_right, content_box)
    
    bottom_box = _box(inner_x, inner_y + top_h + gap, inner_w, bottom_h)
    bottom = _reflow_inside_content_box(bottom_box, content_box)
    assertInsideContentBox(bottom, content_box)
    
    return {"top_left": top_left, "top_right": top_right, "bottom": bottom}


def _add_textbox_clamped(slide, *, box: dict[str, float], content_box: dict[str, float], symbols: dict[str, Any]):
    clamped = _reflow_inside_content_box(box, content_box)
    assertInsideContentBox(clamped, content_box)
    return slide.shapes.add_textbox(
        symbols["Inches"](clamped["x"]),
        symbols["Inches"](clamped["y"]),
        symbols["Inches"](clamped["w"]),
        symbols["Inches"](clamped["h"]),
    )


def _add_picture_clamped(slide, image: BytesIO, *, box: dict[str, float], content_box: dict[str, float], symbols: dict[str, Any]) -> None:
    clamped = _reflow_inside_content_box(box, content_box)
    assertInsideContentBox(clamped, content_box)
    slide.shapes.add_picture(
        image,
        symbols["Inches"](clamped["x"]),
        symbols["Inches"](clamped["y"]),
        width=symbols["Inches"](clamped["w"]),
        height=symbols["Inches"](clamped["h"]),
    )


def _short_header(value: str) -> str:
    mapping = {
        "weekly highlight": "Highlight",
        "delta vs prev": "Delta",
        "building family": "Family",
        "stage category": "Stage",
        "open backlog": "Backlog",
        "approval pct": "Approval",
        "released this week": "Released",
        "released weight this week (t)": "Weight t",
        "piece signal coverage": "Signal %",
        "indexed weight total (kg)": "Weight kg",
        "max age weeks": "Age w",
        "exception type": "Type",
    }
    return mapping.get(str(value).strip().lower(), str(value))


def _fit_table_data(
    headers: list[str],
    rows: list[list[str]],
    *,
    max_cols: int,
    width: float,
) -> tuple[list[str], list[list[str]], list[float]]:
    cols = min(max_cols, len(headers))
    short_headers = [_short_header(h)[:14] for h in headers[:cols]]
    fitted_rows = [[str(col)[:42] for col in row[:cols]] for row in rows]

    ratios: list[float] = []
    for idx in range(cols):
        header_len = max(4, len(short_headers[idx]))
        col_len = max([header_len] + [len(r[idx]) if idx < len(r) else 0 for r in fitted_rows])
        ratios.append(float(max(1, col_len)))

    ratio_sum = sum(ratios) or 1.0
    min_col = max(0.55, width / (cols * 2.5))
    computed = [max(min_col, width * (r / ratio_sum)) for r in ratios]
    scale = width / (sum(computed) or 1.0)
    col_widths = [w * scale for w in computed]
    return short_headers, fitted_rows, col_widths


def _load_pptx_symbols() -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:  # pragma: no cover - environment dependent
        raise RuntimeError(
            "python-pptx is required for executive pack export. Install dependencies from requirements.txt"
        ) from exc

    return {
        "Presentation": Presentation,
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
    }


def _load_matplotlib_plt():
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ModuleNotFoundError as exc:  # pragma: no cover - environment dependent
        raise RuntimeError(
            "matplotlib is required for executive pack chart rendering. Install dependencies from requirements.txt"
        ) from exc
    return plt


def _rgb(rgb_color, triplet: tuple[int, int, int]):
    return rgb_color(*triplet)


def _fmt_pct(value: Any, ndigits: int = 1) -> str:
    try:
        return f"{float(value):.{ndigits}f}%"
    except Exception:
        return "-"


def _fmt_num(value: Any, ndigits: int = 0) -> str:
    try:
        if ndigits <= 0:
            return f"{int(float(value)):,}"
        return f"{float(value):,.{ndigits}f}"
    except Exception:
        return "-"


def _safe_list(value: Any) -> list[dict[str, Any]]:
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    return []


def _severity_rank(value: Any) -> int:
    order = {"critical": 4, "high": 3, "medium": 2, "low": 1}
    return order.get(str(value or "").strip().lower(), 0)


def _render_release_chart(weekly_payload: dict[str, Any]) -> BytesIO:
    plt = _load_matplotlib_plt()
    series = _safe_list(weekly_payload.get("weekly_comparison", {}).get("release_series", []))
    cumulative = _safe_list(weekly_payload.get("weekly_comparison", {}).get("cumulative_series", []))

    weeks = [int(item.get("week")) for item in series if item.get("week") is not None]
    released = [float(item.get("released_dossiers", 0.0) or 0.0) for item in series if item.get("week") is not None]
    cumulative_counts = [
        float(item.get("cumulative_approved_dossiers", 0.0) or 0.0)
        for item in cumulative
        if item.get("week") is not None
    ]
    if len(cumulative_counts) < len(weeks):
        cumulative_counts = cumulative_counts + [0.0] * (len(weeks) - len(cumulative_counts))

    fig, ax = plt.subplots(figsize=(8.2, 3.2), dpi=180)
    if weeks:
        ax.bar(weeks, released, color="#0b2d4e", width=0.75, label="Released dossiers")
        ax2 = ax.twinx()
        ax2.plot(weeks, cumulative_counts[: len(weeks)], color="#008577", linewidth=2.4, label="Cumulative approved")
        ax.set_xlabel("Week")
        ax.set_ylabel("Released")
        ax2.set_ylabel("Cumulative")
        ax.grid(axis="y", alpha=0.24)
    else:
        ax.text(0.5, 0.5, "No weekly release data", ha="center", va="center", fontsize=11)
        ax.set_axis_off()

    fig.tight_layout()
    image = BytesIO()
    fig.savefig(image, format="png", bbox_inches="tight")
    plt.close(fig)
    image.seek(0)
    return image


def _render_signal_alignment_chart(piece_payload: dict[str, Any]) -> BytesIO:
    plt = _load_matplotlib_plt()

    comparison_df = piece_payload.get("comparison")
    if not isinstance(comparison_df, pd.DataFrame):
        comparison_df = pd.DataFrame()

    fig, ax = plt.subplots(figsize=(8.2, 3.2), dpi=180)
    if comparison_df.empty:
        ax.text(0.5, 0.5, "No physical signal comparison data", ha="center", va="center", fontsize=11)
        ax.set_axis_off()
    else:
        work = comparison_df.copy()
        work["gap"] = (work.get("physical_signal_pct", 0.0) - work.get("documented_progress_pct", 0.0)).abs()
        work = work.sort_values("gap", ascending=False).head(8)

        labels = work.get("block", pd.Series(index=work.index, dtype="object")).astype(str).tolist()
        physical = (pd.to_numeric(work.get("physical_signal_pct"), errors="coerce").fillna(0.0) * 100.0).tolist()
        documentary = (
            pd.to_numeric(work.get("documented_progress_pct"), errors="coerce").fillna(0.0) * 100.0
        ).tolist()

        x = list(range(len(labels)))
        width = 0.38
        ax.bar([v - width / 2 for v in x], documentary, width=width, color="#0b2d4e", label="Documented")
        ax.bar([v + width / 2 for v in x], physical, width=width, color="#008577", label="Physical signal")
        ax.set_xticks(x)
        ax.set_xticklabels(labels, rotation=20, ha="right", fontsize=8)
        ax.set_ylim(0, 105)
        ax.set_ylabel("Progress %")
        ax.grid(axis="y", alpha=0.24)
        ax.legend(loc="upper left", fontsize=8)

    fig.tight_layout()
    image = BytesIO()
    fig.savefig(image, format="png", bbox_inches="tight")
    plt.close(fig)
    image.seek(0)
    return image


def _add_brand(slide, *, logo_path: Optional[Path], symbols: dict[str, Any], theme: _Theme) -> None:
    Inches = symbols["Inches"]
    Pt = symbols["Pt"]
    RGBColor = symbols["RGBColor"]

    bar = slide.shapes.add_shape(symbols["MSO_SHAPE"].RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(RGBColor, theme.navy)
    bar.line.fill.background()

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.55), Inches(0.02), height=Inches(0.26))

    label = slide.shapes.add_textbox(Inches(0.35), Inches(0.04), Inches(6.5), Inches(0.24))
    tf = label.text_frame
    tf.text = "INPROS | QA Platform"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(RGBColor, theme.white)


def _add_title(slide, title: str, subtitle: str, *, symbols: dict[str, Any], theme: _Theme) -> None:
    Inches = symbols["Inches"]
    Pt = symbols["Pt"]
    RGBColor = symbols["RGBColor"]

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(9.6), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = _rgb(RGBColor, theme.navy)

    subtitle_box = slide.shapes.add_textbox(Inches(CONTENT_BOX["x"]), Inches(1.18), Inches(CONTENT_BOX["w"]), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(12)
    subtitle_frame.paragraphs[0].font.color.rgb = _rgb(RGBColor, theme.slate)


def _add_kpi_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    value: str,
    content_box: dict[str, float],
    symbols: dict[str, Any],
    theme: _Theme,
) -> None:
    Inches = symbols["Inches"]
    Pt = symbols["Pt"]
    RGBColor = symbols["RGBColor"]

    card_box = _reflow_inside_content_box(_box(left, top, width, height), content_box)
    card = slide.shapes.add_shape(
        symbols["MSO_SHAPE"].ROUNDED_RECTANGLE,
        Inches(card_box["x"]),
        Inches(card_box["y"]),
        Inches(card_box["w"]),
        Inches(card_box["h"]),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(RGBColor, theme.light_bg)
    card.line.color.rgb = _rgb(RGBColor, (214, 220, 227))
    assertInsideContentBox(card_box, content_box)

    label_box = _add_textbox_clamped(
        slide,
        box=_box(card_box["x"] + 0.12, card_box["y"] + 0.08, max(0.2, card_box["w"] - 0.18), 0.20),
        content_box=content_box,
        symbols=symbols,
    )
    label_tf = label_box.text_frame
    label_tf.text = label
    label_tf.paragraphs[0].font.size = Pt(9)
    label_tf.paragraphs[0].font.bold = True
    label_tf.paragraphs[0].font.color.rgb = _rgb(RGBColor, theme.slate)

    value_box = _add_textbox_clamped(
        slide,
        box=_box(card_box["x"] + 0.12, card_box["y"] + 0.30, max(0.2, card_box["w"] - 0.18), 0.38),
        content_box=content_box,
        symbols=symbols,
    )
    value_tf = value_box.text_frame
    value_tf.text = value
    value_tf.paragraphs[0].font.size = Pt(18)
    value_tf.paragraphs[0].font.bold = True
    value_tf.paragraphs[0].font.color.rgb = _rgb(RGBColor, theme.navy)


def _add_table(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    headers: list[str],
    rows: list[list[str]],
    content_box: dict[str, float],
    max_cols: Optional[int],
    symbols: dict[str, Any],
    theme: _Theme,
) -> None:
    Inches = symbols["Inches"]
    Pt = symbols["Pt"]
    RGBColor = symbols["RGBColor"]

    table_box = _reflow_inside_content_box(_box(left, top, width, height), content_box)
    assertInsideContentBox(table_box, content_box)

    headers_fit, rows_fit, col_widths = _fit_table_data(
        headers,
        rows,
        max_cols=max_cols or len(headers),
        width=table_box["w"],
    )

    table_shape = slide.shapes.add_table(
        len(rows_fit) + 1,
        len(headers_fit),
        Inches(table_box["x"]),
        Inches(table_box["y"]),
        Inches(table_box["w"]),
        Inches(table_box["h"]),
    )
    table = table_shape.table

    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = Inches(col_width)

    for idx, header in enumerate(headers_fit):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(RGBColor, theme.navy)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = _rgb(RGBColor, theme.white)
        cell.margin_left = Inches(0.02)
        cell.margin_right = Inches(0.02)
        cell.margin_top = Inches(0.01)
        cell.margin_bottom = Inches(0.01)

    for ridx, row in enumerate(rows_fit, start=1):
        for cidx, value in enumerate(row):
            cell = table.cell(ridx, cidx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(RGBColor, theme.white if ridx % 2 else theme.light_bg)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(8)
            run.font.color.rgb = _rgb(RGBColor, theme.navy)
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)


def _add_cover_slide(prs, *, report_meta: dict[str, Any], symbols: dict[str, Any], theme: _Theme, logo_path: Optional[Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand(slide, logo_path=logo_path, symbols=symbols, theme=theme)

    analysis_week = report_meta.get("analysis_week")
    generated_at = report_meta.get("generated_at") or datetime.now(timezone.utc).isoformat()
    subtitle = f"Board Executive Pack | Week {analysis_week if analysis_week is not None else '-'} | Generated {generated_at[:19]}Z"
    _add_title(slide, "QA Platform Executive Pack", subtitle, symbols=symbols, theme=theme)

    body = _add_textbox_clamped(
        slide,
        box=_box(CONTENT_BOX["x"], 2.0, CONTENT_BOX["w"], 1.8),
        content_box=CONTENT_BOX,
        symbols=symbols,
    )
    tf = body.text_frame
    tf.text = "Executive narrative: performance, weekly execution signal, and material risks."
    tf.paragraphs[0].font.size = symbols["Pt"](15)
    tf.paragraphs[0].font.color.rgb = _rgb(symbols["RGBColor"], theme.slate)


def _add_kpi_slide(
    prs,
    *,
    executive_payload: dict[str, Any],
    symbols: dict[str, Any],
    theme: _Theme,
    logo_path: Optional[Path],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand(slide, logo_path=logo_path, symbols=symbols, theme=theme)
    _add_title(slide, "Executive KPI Overview", "Contract scope performance at a glance", symbols=symbols, theme=theme)

    kpis = executive_payload.get("executive_kpis", {})
    weekly = executive_payload.get("weekly_management", {})

    cards = [
        ("Total Dossiers", _fmt_num(kpis.get("total_dossiers", 0))),
        ("Approved", _fmt_num(kpis.get("approved_dossiers", 0))),
        ("Open Backlog", _fmt_num(weekly.get("backlog_aging_summary", {}).get("total_open_backlog", 0))),
        ("Weight Released (t)", _fmt_num(kpis.get("peso_liberado_ton", 0.0), ndigits=2)),
    ]

    overview_layout = _layout_overview(CONTENT_BOX)
    for idx, (label, value) in enumerate(cards):
        card_box = overview_layout["cards"][idx]
        _add_kpi_card(
            slide,
            left=card_box["x"],
            top=card_box["y"],
            width=card_box["w"],
            height=card_box["h"],
            label=label,
            value=value,
            content_box=CONTENT_BOX,
            symbols=symbols,
            theme=theme,
        )

    highlights = _safe_list(executive_payload.get("weekly_highlights", []))
    rows = [
        [
            str(item.get("label", "-")).replace("_", " ").title(),
            _fmt_num(item.get("value", 0), ndigits=2 if "weight" in str(item.get("label", "")) else 0),
            _fmt_num(item.get("delta", 0), ndigits=2 if "weight" in str(item.get("label", "")) else 0),
        ]
        for item in highlights[:6]
    ]
    if not rows:
        rows = [["No highlights", "-", "-"]]

    _add_table(
        slide,
        left=overview_layout["left_table"]["x"],
        top=overview_layout["left_table"]["y"],
        width=overview_layout["left_table"]["w"],
        height=overview_layout["left_table"]["h"],
        headers=["Weekly Highlight", "Value", "Delta vs Prev"],
        rows=rows,
        content_box=CONTENT_BOX,
        max_cols=3,
        symbols=symbols,
        theme=theme,
    )

    summary_rows = []
    for row in _safe_list(executive_payload.get("executive_summary_table", []))[:6]:
        summary_rows.append(
            [
                str(row.get("stage_category", "-")),
                str(row.get("building_family", "-")),
                _fmt_num(row.get("open_backlog", 0)),
                _fmt_pct(row.get("approval_pct", 0.0)),
            ]
        )
    if not summary_rows:
        summary_rows = [["-", "-", "-", "-"]]

    _add_table(
        slide,
        left=overview_layout["right_table"]["x"],
        top=overview_layout["right_table"]["y"],
        width=overview_layout["right_table"]["w"],
        height=overview_layout["right_table"]["h"],
        headers=["Stage", "Family", "Backlog", "Approval"],
        rows=summary_rows,
        content_box=CONTENT_BOX,
        max_cols=4,
        symbols=symbols,
        theme=theme,
    )


def _add_weekly_signal_slide(
    prs,
    *,
    executive_payload: dict[str, Any],
    piece_payload: dict[str, Any],
    symbols: dict[str, Any],
    theme: _Theme,
    logo_path: Optional[Path],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand(slide, logo_path=logo_path, symbols=symbols, theme=theme)
    _add_title(slide, "Weekly Management + Physical Execution Signal", "Documentary trend against field signal", symbols=symbols, theme=theme)

    release_chart = _render_release_chart(executive_payload.get("weekly_management", {}))
    signal_chart = _render_signal_alignment_chart(piece_payload)

    include_secondary_chart = True
    trend_layout = _layout_trend(CONTENT_BOX, include_secondary_chart=include_secondary_chart)

    main_chart_box = trend_layout["main_chart"]
    secondary_chart_box = trend_layout["secondary_chart"]
    assertInsideContentBox(main_chart_box, CONTENT_BOX)
    assertInsideContentBox(secondary_chart_box, CONTENT_BOX)

    _add_picture_clamped(
        slide,
        release_chart,
        box=main_chart_box,
        content_box=CONTENT_BOX,
        symbols=symbols,
    )
    _add_picture_clamped(
        slide,
        signal_chart,
        box=secondary_chart_box,
        content_box=CONTENT_BOX,
        symbols=symbols,
    )

    kpis = piece_payload.get("kpis", {}) if isinstance(piece_payload.get("kpis"), dict) else {}
    weekly = executive_payload.get("weekly_management", {})
    delta = weekly.get("delta_kpis", {}) if isinstance(weekly.get("delta_kpis"), dict) else {}

    metric_rows = [
        ["Released this week", _fmt_num(delta.get("released_this_week", 0))],
        ["Released weight this week (t)", _fmt_num(delta.get("released_weight_t_this_week", 0.0), ndigits=2)],
        ["Piece signal coverage", _fmt_pct(float(kpis.get("week_trace_coverage_pct", 0.0) or 0.0) * 100.0, ndigits=1)],
        ["Indexed weight total (kg)", _fmt_num(kpis.get("indexed_weight_total", 0.0), ndigits=0)],
    ]

    _add_table(
        slide,
        left=trend_layout["bottom_table"]["x"],
        top=trend_layout["bottom_table"]["y"],
        width=trend_layout["bottom_table"]["w"],
        height=trend_layout["bottom_table"]["h"],
        headers=["Weekly Metric", "Value"],
        rows=metric_rows,
        content_box=CONTENT_BOX,
        max_cols=2,
        symbols=symbols,
        theme=theme,
    )


def _add_risk_slide(
    prs,
    *,
    executive_payload: dict[str, Any],
    piece_payload: dict[str, Any],
    symbols: dict[str, Any],
    theme: _Theme,
    logo_path: Optional[Path],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand(slide, logo_path=logo_path, symbols=symbols, theme=theme)
    _add_title(slide, "Risk and Exceptions", "Top-N concentration for board focus", symbols=symbols, theme=theme)

    top_backlog = _safe_list(executive_payload.get("top_backlog_risks", []))[:5]
    backlog_rows = [
        [
            str(item.get("stage_category", "-")),
            str(item.get("building_family", "-")),
            _fmt_num(item.get("open_backlog", 0)),
            _fmt_num(item.get("max_age_weeks", 0)),
        ]
        for item in top_backlog
    ]
    if not backlog_rows:
        backlog_rows = [["-", "-", "-", "-"]]

    approval_gaps = _safe_list(executive_payload.get("risk_exception_summary", {}).get("largest_approval_gaps", []))[:5]
    approval_rows = [
        [
            str(item.get("stage_category", "-")),
            str(item.get("building_family", "-")),
            _fmt_num(item.get("approval_gap", 0)),
            _fmt_pct(item.get("approval_pct", 0.0)),
        ]
        for item in approval_gaps
    ]
    if not approval_rows:
        approval_rows = [["-", "-", "-", "-"]]

    exceptions_df = piece_payload.get("exceptions")
    if not isinstance(exceptions_df, pd.DataFrame):
        exceptions_df = pd.DataFrame(columns=["severity", "exception_type", "block", "details"])
    exception_rows_df = exceptions_df.copy()
    if not exception_rows_df.empty:
        exception_rows_df["_rank"] = exception_rows_df.get("severity", "").apply(_severity_rank)
        exception_rows_df = exception_rows_df.sort_values(["_rank", "exception_type"], ascending=[False, True])
    exception_rows = [
        [
            str(row.get("severity", "-")).upper(),
            str(row.get("exception_type", "-"))[:38],
            str(row.get("block", "-")),
            str(row.get("details", "-"))[:46],
        ]
        for _, row in exception_rows_df.head(6).iterrows()
    ]
    if not exception_rows:
        exception_rows = [["-", "-", "-", "-"]]

    risk_layout = _layout_risk(CONTENT_BOX)

    _add_table(
        slide,
        left=risk_layout["top_left"]["x"],
        top=risk_layout["top_left"]["y"],
        width=risk_layout["top_left"]["w"],
        height=risk_layout["top_left"]["h"],
        headers=["Stage", "Family", "Backlog", "Age w"],
        rows=backlog_rows,
        content_box=CONTENT_BOX,
        max_cols=4,
        symbols=symbols,
        theme=theme,
    )

    _add_table(
        slide,
        left=risk_layout["top_right"]["x"],
        top=risk_layout["top_right"]["y"],
        width=risk_layout["top_right"]["w"],
        height=risk_layout["top_right"]["h"],
        headers=["Stage", "Family", "Gap", "Approval"],
        rows=approval_rows,
        content_box=CONTENT_BOX,
        max_cols=4,
        symbols=symbols,
        theme=theme,
    )

    _add_table(
        slide,
        left=risk_layout["bottom"]["x"],
        top=risk_layout["bottom"]["y"],
        width=risk_layout["bottom"]["w"],
        height=risk_layout["bottom"]["h"],
        headers=["Severity", "Type", "Block", "Details"],
        rows=exception_rows,
        content_box=CONTENT_BOX,
        max_cols=4,
        symbols=symbols,
        theme=theme,
    )


def _add_closing_slide(
    prs,
    *,
    executive_payload: dict[str, Any],
    symbols: dict[str, Any],
    theme: _Theme,
    logo_path: Optional[Path],
) -> None:
    insights = _safe_list(executive_payload.get("high_value_insights", []))
    risk_signals = _safe_list(executive_payload.get("risk_exception_summary", {}).get("signals", []))
    if not insights and not risk_signals:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand(slide, logo_path=logo_path, symbols=symbols, theme=theme)
    _add_title(slide, "Alignment and Closing Summary", "Action-oriented executive takeaways", symbols=symbols, theme=theme)

    lines: list[str] = []
    for item in insights[:4]:
        key = str(item.get("key", "insight")).replace("_", " ").title()
        value = item.get("value")
        context_bits = []
        for field in ("stage_category", "building_family", "context"):
            raw = item.get(field)
            if raw not in (None, ""):
                context_bits.append(str(raw))
        suffix = f" ({', '.join(context_bits)})" if context_bits else ""
        lines.append(f"- {key}: {value}{suffix}")

    for item in risk_signals[:3]:
        key = str(item.get("key", "signal")).replace("_", " ").title()
        value = item.get("value", "-")
        context = item.get("context", "")
        context_txt = f" {context}" if context else ""
        lines.append(f"- Risk signal: {key} = {value}{context_txt}")

    if not lines:
        return

    box = _add_textbox_clamped(
        slide,
        box=_box(CONTENT_BOX["x"], 1.9, CONTENT_BOX["w"], 4.8),
        content_box=CONTENT_BOX,
        symbols=symbols,
    )
    tf = box.text_frame
    tf.text = lines[0]
    tf.paragraphs[0].font.size = symbols["Pt"](15)
    tf.paragraphs[0].font.color.rgb = _rgb(symbols["RGBColor"], theme.navy)

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = symbols["Pt"](13)
        p.font.color.rgb = _rgb(symbols["RGBColor"], theme.slate)


def generate_executive_pack_pptx(
    *,
    week: Optional[int] = None,
    comparison_week: Optional[int] = None,
    language: str = "en",
) -> bytes:
    symbols = _load_pptx_symbols()
    theme = _Theme()

    executive_payload = build_executive_report_payload(
        selected_week=week,
        comparison_week=comparison_week,
        language=language,
    )

    try:
        piece_payload = load_piece_signal_payload(rebuild_if_missing=False)
    except Exception:
        piece_payload = {
            "kpis": {},
            "comparison": pd.DataFrame(),
            "exceptions": pd.DataFrame(),
        }

    Presentation = symbols["Presentation"]
    prs = Presentation()

    settings = get_settings()
    logo_path = settings.project_root / "assets" / "inpros-logo.png"

    _add_cover_slide(prs, report_meta=executive_payload.get("report_meta", {}), symbols=symbols, theme=theme, logo_path=logo_path)
    _add_kpi_slide(prs, executive_payload=executive_payload, symbols=symbols, theme=theme, logo_path=logo_path)
    _add_weekly_signal_slide(
        prs,
        executive_payload=executive_payload,
        piece_payload=piece_payload,
        symbols=symbols,
        theme=theme,
        logo_path=logo_path,
    )
    _add_risk_slide(
        prs,
        executive_payload=executive_payload,
        piece_payload=piece_payload,
        symbols=symbols,
        theme=theme,
        logo_path=logo_path,
    )
    _add_closing_slide(prs, executive_payload=executive_payload, symbols=symbols, theme=theme, logo_path=logo_path)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()


def executive_pack_filename(analysis_week: Any) -> str:
    week_value = str(analysis_week) if analysis_week not in (None, "") else "NA"
    return f"executive-pack-week-{week_value}.pptx"
